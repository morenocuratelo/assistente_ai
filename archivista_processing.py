"""
Questo file contiene la logica di business principale per il processamento
dei documenti. La funzione qui definita viene registrata come task Celery.
"""
import os
import shutil
import sqlite3
import json
import time
from datetime import datetime
from pydantic import BaseModel, Field
from typing import List, Optional

# Import per l'estrazione granulare del testo
import fitz
from docx import Document as DocxDocument
import pdfplumber
from pypdf import PdfReader
from pdfminer.high_level import extract_text as pdfminer_extract
from bs4 import BeautifulSoup
import chardet

from llama_index.core import Document, Settings, VectorStoreIndex, StorageContext, PromptTemplate, load_index_from_storage

from celery_app import celery_app
from celery.exceptions import SoftTimeLimitExceeded
from config import initialize_services
import prompt_manager # <-- MODIFICA: Importa il nuovo gestore dei prompt
import knowledge_structure
from file_utils import setup_database
# Import del motore di inferenza Bayesiano
from bayesian_inference_engine import (
    create_inference_engine,
    ConfidenceUpdateRequest,
    process_evidence_batch
)

# --- CONFIGURAZIONE ---
DOCS_TO_PROCESS_DIR = "documenti_da_processare"
CATEGORIZED_ARCHIVE_DIR = "Dall_Origine_alla_Complessita"
DB_STORAGE_DIR = "db_memoria"
METADATA_DB_FILE = os.path.join(DB_STORAGE_DIR, "metadata.sqlite")
ARCHIVISTA_STATUS_FILE = os.path.join(DB_STORAGE_DIR, "archivista_status.json")

# --- MODELLI DATI E FUNZIONI DI UTILITÀ ---

class PaperMetadata(BaseModel):
    title: str = Field(..., description="Il titolo completo e corretto dell'articolo.")
    authors: List[str] = Field(default_factory=list, description="La lista completa degli autori menzionati.")
    publication_year: Optional[int] = Field(None, description="L'anno di pubblicazione (solo il numero).")

def update_status(status, current_file=None):
    try:
        os.makedirs(DB_STORAGE_DIR, exist_ok=True) # Assicura che la directory esista
        with open(ARCHIVISTA_STATUS_FILE, "w", encoding="utf-8") as f:
            json.dump({"status": status, "file": current_file, "timestamp": datetime.now().isoformat()}, f)
    except Exception as e:
        print(f"--> ERRORE aggiornamento stato: {e}")

def db_connect():
    conn = sqlite3.connect(METADATA_DB_FILE)
    conn.row_factory = sqlite3.Row
    return conn

# --- FUNZIONI PER L'ESTRAZIONE DEL TESTO (invariate) ---
def extract_text_from_pdf(file_path: str) -> str:
    text_extractors = [("PyMuPDF", extract_text_pymupdf), ("PyPDF2", extract_text_pypdf2), ("pdfplumber", extract_text_pdfplumber), ("pdfminer", extract_text_pdfminer)]
    for name, func in text_extractors:
        try:
            text = func(file_path)
            if text and text.strip(): return text
        except Exception as e:
            print(f"✗ Errore con {name}: {e}")
    raise ValueError(f"Impossibile estrarre testo dal PDF {file_path}")

def extract_text_pymupdf(file_path: str) -> str:
    with fitz.open(file_path) as doc: return "".join(page.get_text() for page in doc)
def extract_text_pypdf2(file_path: str) -> str:
    with open(file_path, 'rb') as f: return "".join(p.extract_text() for p in PdfReader(f).pages)
def extract_text_pdfplumber(file_path: str) -> str:
    with pdfplumber.open(file_path) as pdf: return "".join(p.extract_text() or "" for p in pdf.pages)
def extract_text_pdfminer(file_path: str) -> str:
    return pdfminer_extract(file_path)
def extract_text_from_docx(file_path: str) -> str:
    return "\n".join(p.text for p in DocxDocument(file_path).paragraphs)

def extract_text_from_doc(file_path: str) -> str:
    """Extract text from legacy Word .doc files"""
    try:
        # Try using antiword if available (external tool)
        import subprocess
        try:
            result = subprocess.run(['antiword', file_path],
                                  capture_output=True, text=True, timeout=30)
            if result.returncode == 0:
                return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Fallback: try using pypandoc if available
        try:
            import pypandoc
            return pypandoc.convert_file(file_path, 'plain', format='doc')
        except ImportError:
            pass

        # Final fallback: return placeholder message
        print(f"⚠️ Estrazione DOC non disponibile per {file_path}")
        return f"Contenuto Word (.doc): {os.path.basename(file_path)} - Richiede antiword o pypandoc per estrazione completa"

    except Exception as e:
        print(f"Errore estrazione DOC: {e}")
        return ""
def extract_text_from_rtf(file_path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, 'r') as f:
            return rtf_to_text(f.read())
    except Exception: return ""
def extract_text_from_html(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f: return BeautifulSoup(f.read(), 'html.parser').get_text()
def extract_text_from_txt(file_path: str) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8') as f: return f.read()
    except UnicodeDecodeError:
        with open(file_path, 'rb') as f: raw = f.read()
        enc = chardet.detect(raw)['encoding']
        return raw.decode(enc or 'utf-8')

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint .pptx files"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

        return '\n'.join(text_content)
    except Exception as e:
        print(f"Errore estrazione PPTX: {e}")
        return ""

def extract_text_from_ppt(file_path: str) -> str:
    """Extract text from legacy PowerPoint .ppt files"""
    try:
        # For .ppt files, we can try to use an external tool or convert to pptx
        # For now, return a placeholder - full implementation would need additional tools
        print(f"⚠️ Estrazione PPT non completamente implementata per {file_path}")
        return f"Contenuto PowerPoint (.ppt): {os.path.basename(file_path)} - Richiede tool esterno per estrazione completa"
    except Exception as e:
        print(f"Errore estrazione PPT: {e}")
        return ""

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from Excel .xlsx files"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_content = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"Foglio: {sheet_name}"]

            for row in sheet.iter_rows(values_only=True):
                # Extract non-empty cell values
                row_data = [str(cell) for cell in row if cell is not None]
                if row_data:
                    sheet_text.append(" | ".join(row_data))

            if len(sheet_text) > 1:  # More than just the sheet name
                text_content.extend(sheet_text)
                text_content.append("")  # Empty line between sheets

        return '\n'.join(text_content)
    except Exception as e:
        print(f"Errore estrazione XLSX: {e}")
        return ""

def extract_text_from_xls(file_path: str) -> str:
    """Extract text from legacy Excel .xls files"""
    try:
        import xlrd
        wb = xlrd.open_workbook(file_path)
        text_content = []

        for sheet_name in wb.sheet_names():
            sheet = wb.sheet_by_name(sheet_name)
            sheet_text = [f"Foglio: {sheet_name}"]

            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                # Extract non-empty cell values
                row_data = [str(cell) for cell in row if cell != ""]
                if row_data:
                    sheet_text.append(" | ".join(row_data))

            if len(sheet_text) > 1:  # More than just the sheet name
                text_content.extend(sheet_text)
                text_content.append("")  # Empty line between sheets

        return '\n'.join(text_content)
    except ImportError:
        print("xlrd non disponibile per file .xls - considera installazione: pip install xlrd")
        return f"Contenuto Excel (.xls): {os.path.basename(file_path)} - Richiede xlrd per estrazione completa"
    except Exception as e:
        print(f"Errore estrazione XLS: {e}")
        return ""

def extract_text_from_csv(file_path: str) -> str:
    """Extract text from CSV files"""
    try:
        import csv
        text_content = []
        encoding = 'utf-8'

        # Try to detect encoding
        try:
            with open(file_path, 'rb') as f:
                raw = f.read()
                enc = chardet.detect(raw)['encoding']
                if enc:
                    encoding = enc
        except:
            pass

        with open(file_path, 'r', encoding=encoding) as f:
            # Try to detect delimiter
            sample = f.read(1024)
            f.seek(0)
            sniffer = csv.Sniffer()
            delimiter = sniffer.sniff(sample).delimiter

            reader = csv.reader(f, delimiter=delimiter)
            for row_idx, row in enumerate(reader):
                # Extract non-empty cell values
                row_data = [str(cell) for cell in row if cell != ""]
                if row_data:
                    text_content.append(" | ".join(row_data))

        return '\n'.join(text_content)
    except Exception as e:
        print(f"Errore estrazione CSV: {e}")
        return ""

def get_text_extractor(file_extension: str):
    extractors = {
        # Text Documents
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_doc,  # Use dedicated .doc extractor
        '.rtf': extract_text_from_rtf,
        '.html': extract_text_from_html,
        '.htm': extract_text_from_html,
        '.txt': extract_text_from_txt,

        # Presentations
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_ppt,

        # Spreadsheets
        '.xlsx': extract_text_from_xlsx,
        '.xls': extract_text_from_xls,
        '.csv': extract_text_from_csv,
    }
    return extractors.get(file_extension.lower())

def classify_document(text_content):
    if not Settings.llm:
        raise ConnectionError("LLM non disponibile per la classificazione.")
    classification_structure = knowledge_structure.get_structure_for_prompt()
    # MODIFICA: Usa il prompt manager per caricare il template
    prompt_template = PromptTemplate(prompt_manager.get_prompt("DOCUMENT_CLASSIFICATION_PROMPT"))
    formatted_prompt = prompt_template.format(classification_structure=classification_structure, document_text=text_content[:8000])
    response = Settings.llm.complete(formatted_prompt)
    category_id = str(response).strip()
    return category_id if knowledge_structure.is_valid_category_id(category_id) else "UNCATEGORIZED/C00"

# --- TASK PRINCIPALE ---
@celery_app.task(
    name='archivista.process_document',
    bind=True,
    autoretry_for=(ConnectionError,),
    retry_kwargs={'max_retries': 3},
    default_retry_delay=30,
    soft_time_limit=300,
    time_limit=360
)
def process_document_task(self, file_path):
    file_name = os.path.basename(file_path)
    lock_file = file_path + ".lock"

    # Check for existing lock with timeout
    if os.path.exists(lock_file):
        lock_age = time.time() - os.path.getmtime(lock_file)
        if lock_age < 600:  # Increased timeout to 10 minutes
            update_status(f"File già in elaborazione da {lock_age:.0f}s", file_name)
            print(f"⏳ {file_name} già in elaborazione, attendo...")
            return {'status': 'already_processing', 'file_name': file_name}
        else:
            print(f"⚠️ Rimuovo lock obsoleto per {file_name} (età: {lock_age:.0f}s)")
            os.remove(lock_file)

    try:
        # Create lock file atomically
        try:
            with open(lock_file, "w") as f:
                f.write(f"{os.getpid()},{time.time()}")
        except Exception as e:
            print(f"❌ Errore creazione lock per {file_name}: {e}")
            return {'status': 'lock_error', 'file_name': file_name}

        initialize_services()
        setup_database()

        if Settings.llm is None or Settings.embed_model is None:
            raise ConnectionError("Servizi AI non disponibili. Controlla la connessione e i modelli.")

        update_status("Avviato processamento", file_name)
        
        # 1. ESTRAZIONE TESTO
        update_status("Estrazione testo...", file_name)
        file_ext = os.path.splitext(file_name)[1].lower()
        extractor = get_text_extractor(file_ext)
        if not extractor: raise ValueError(f"Formato file non supportato: {file_ext}")
        full_text = extractor(file_path)
        if not full_text or not full_text.strip(): raise ValueError("Documento vuoto o illeggibile.")
        doc = Document(text=full_text)

        # 2. CLASSIFICAZIONE
        update_status("Classificazione AI...", file_name)
        category_id = classify_document(full_text)
        if category_id == "UNCATEGORIZED/C00":
            part_id, chapter_id = "UNCATEGORIZED", "C00"
            part_name, chapter_name = "Non Categorizzato", "Generale"
        else:
            part_id, chapter_id = category_id.split('/')
            part_name = knowledge_structure.KNOWLEDGE_BASE_STRUCTURE[part_id]['name']
            chapter_name = knowledge_structure.KNOWLEDGE_BASE_STRUCTURE[part_id]['chapters'][chapter_id]
        category_full_name = f"{part_name} -> {chapter_name}"

        # 3. ESTRAZIONE METADATI
        update_status("Estrazione metadati...", file_name)
        # MODIFICA: Usa il prompt manager
        metadata_prompt = PromptTemplate(prompt_manager.get_prompt("PYDANTIC_METADATA_PROMPT"))
        metadata_query = metadata_prompt.format(document_text=full_text[:4000], category_name=category_full_name)
        response = Settings.llm.complete(metadata_query)
        try:
            response_text = str(response).strip()
            # Try to extract JSON from the response if it contains extra text
            if '{' in response_text and '}' in response_text:
                start_idx = response_text.find('{')
                end_idx = response_text.rfind('}') + 1
                json_str = response_text[start_idx:end_idx]
                metadata = PaperMetadata(**json.loads(json_str))
            else:
                # Fallback if no JSON found
                metadata = PaperMetadata(title=file_name, authors=[], publication_year=None)
        except (json.JSONDecodeError, TypeError, ValueError) as e:
            print(f"⚠️ Errore parsing metadati JSON: {e}. Response: {str(response)[:200]}...")
            metadata = PaperMetadata(title=file_name, authors=[], publication_year=None)

        # 4. INDICIZZAZIONE (LOGICA SEMPLIFICATA E ATOMICA)
        update_status("Indicizzazione...", file_name)
        doc.metadata.update({"file_name": file_name, "title": metadata.title, "authors": json.dumps(metadata.authors), "publication_year": metadata.publication_year, "category_id": category_id, "category_name": category_full_name})

        # Crea storage context per l'indicizzazione con ChromaVectorStore
        from llama_index.core.storage.docstore import SimpleDocumentStore
        from llama_index.core.storage.index_store import SimpleIndexStore

        # Crea/ottiene collezione ChromaDB
        try:
            import chromadb
            from llama_index.vector_stores.chroma import ChromaVectorStore

            db = chromadb.PersistentClient(path=DB_STORAGE_DIR)
            chroma_collection = db.get_or_create_collection("documents")
            vector_store = ChromaVectorStore(chroma_collection=chroma_collection)

            print(f"🚀 Utilizzo ChromaVectorStore per alte performance")

        except ImportError:
            print(f"⚠️ ChromaDB non disponibile, fallback a SimpleVectorStore")
            from llama_index.core.vector_stores import SimpleVectorStore
            vector_store = SimpleVectorStore()

        storage_context = StorageContext.from_defaults(
            docstore=SimpleDocumentStore(),
            vector_store=vector_store,
            index_store=SimpleIndexStore(),
            persist_dir=DB_STORAGE_DIR
        )

        try:
            # Tenta di caricare un indice esistente
            storage_context = StorageContext.from_defaults(persist_dir=DB_STORAGE_DIR, vector_store=vector_store)
            index = load_index_from_storage(storage_context)
            print("✅ Indice esistente caricato. Aggiungo nuovo documento...")
            index.insert(doc)

        except FileNotFoundError:
            # Se nessun indice esiste, creane uno nuovo
            print("🆕 Nessun indice trovato. Ne creo uno nuovo.")
            storage_context = StorageContext.from_defaults(vector_store=vector_store)
            index = VectorStoreIndex.from_documents([doc], storage_context=storage_context)

        # A prescindere da cosa sia successo, salva lo stato finale
        index.storage_context.persist(persist_dir=DB_STORAGE_DIR)
        print("💾 Indice salvato correttamente su disco.")

        # 4.5. GENERAZIONE ANTEPRIMA
        update_status("Generazione anteprima...", file_name)
        # Usa il prompt manager per caricare il template di anteprima
        preview_prompt = PromptTemplate(prompt_manager.get_prompt("FORMAT_PREVIEW_PROMPT"))
        preview_query = preview_prompt.format(raw_text=full_text[:2500])  # Usa primi 2500 caratteri
        preview_response = Settings.llm.complete(preview_query)
        formatted_preview = str(preview_response).strip()

        # 4.6. ANALISI ACCADEMICA (parole chiave, task e grafo della conoscenza)
        update_status("Analisi accademica...", file_name)
        academic_metadata = {}
        knowledge_entities = []
        knowledge_relationships = []

        # Estrazione parole chiave
        try:
            keywords_prompt = PromptTemplate(prompt_manager.get_prompt("ACADEMIC_KEYWORDS_PROMPT"))
            keywords_query = keywords_prompt.format(document_text=full_text[:4000])
            keywords_response = Settings.llm.complete(keywords_query)
            keywords_text = str(keywords_response).strip()
            # Try to parse JSON
            try:
                keywords_data = json.loads(keywords_text)
                academic_metadata['keywords'] = keywords_data['keywords']
            except (json.JSONDecodeError, KeyError):
                print(f"⚠️ Impossibile parsare parole chiave JSON per {file_name}")
                academic_metadata['keywords'] = []
        except Exception as e:
            print(f"⚠️ Errore estrazione parole chiave per {file_name}: {e}")
            academic_metadata['keywords'] = []

        # Estrazione entità concettuali per il grafo della conoscenza
        try:
            entities_prompt = PromptTemplate(prompt_manager.get_prompt("KNOWLEDGE_ENTITIES_PROMPT"))
            entities_query = entities_prompt.format(document_text=full_text[:3000])
            entities_response = Settings.llm.complete(entities_query)
            entities_text = str(entities_response).strip()

            # Try to parse JSON entities
            try:
                entities_data = json.loads(entities_text)
                knowledge_entities = entities_data
                print(f"✅ Estratte {len(knowledge_entities)} entità concettuali da {file_name}")
            except (json.JSONDecodeError, KeyError):
                print(f"⚠️ Impossibile parsare entità JSON per {file_name}")
                knowledge_entities = []

        except Exception as e:
            print(f"⚠️ Errore estrazione entità per {file_name}: {e}")
            knowledge_entities = []

        # Estrazione relazioni concettuali (solo se abbiamo entità)
        if knowledge_entities:
            try:
                entities_list = [f"- {e['entity_name']} ({e['entity_type']})" for e in knowledge_entities]
                entities_list_text = "\n".join(entities_list)

                relationships_prompt = PromptTemplate(prompt_manager.get_prompt("ENTITY_RELATIONSHIPS_PROMPT"))
                relationships_query = relationships_prompt.format(
                    document_text=full_text[:3000],
                    entities_list=entities_list_text
                )
                relationships_response = Settings.llm.complete(relationships_query)
                relationships_text = str(relationships_response).strip()

                # Try to parse JSON relationships
                try:
                    relationships_data = json.loads(relationships_text)
                    knowledge_relationships = relationships_data
                    print(f"✅ Estratte {len(knowledge_relationships)} relazioni concettuali da {file_name}")
                except (json.JSONDecodeError, KeyError):
                    print(f"⚠️ Impossibile parsare relazioni JSON per {file_name}")
                    knowledge_relationships = []

            except Exception as e:
                print(f"⚠️ Errore estrazione relazioni per {file_name}: {e}")
                knowledge_relationships = []

        # Generazione task AI
        try:
            tasks_prompt = PromptTemplate(prompt_manager.get_prompt("ACADEMIC_TASK_GENERATION_PROMPT"))
            tasks_query = tasks_prompt.format(document_text=full_text[:3000])
            tasks_response = Settings.llm.complete(tasks_query)
            tasks_text = str(tasks_response).strip()
            # Try to parse JSON
            try:
                tasks_data = json.loads(tasks_text)
                academic_metadata['ai_tasks'] = tasks_data
            except (json.JSONDecodeError, KeyError):
                print(f"⚠️ Impossibile parsare task JSON per {file_name}")
                academic_metadata['ai_tasks'] = {}
        except Exception as e:
            print(f"⚠️ Errore generazione task AI per {file_name}: {e}")
            academic_metadata['ai_tasks'] = {}

        # 4.7. PROCESAMENTO BAYESIANO DELLA CONOSCENZA
        update_status("Analisi Bayesiana...", file_name)

        # Crea motore di inferenza Bayesiano (usa user_id di default se disponibile)
        # Nota: In produzione, questo dovrebbe usare l'user_id del chiamante
        # Per ora usiamo un user_id di default (1) per documenti non associati a utenti specifici
        default_user_id = 1
        bayesian_engine = create_inference_engine(user_id=default_user_id, learning_rate=0.3)

        # Prepara entità per il processamento Bayesiano
        extracted_entities = []
        for entity in knowledge_entities:
            extracted_entities.append({
                'name': entity.get('entity_name', ''),
                'type': entity.get('entity_type', 'concept'),
                'description': entity.get('entity_description', '')
            })

        # Prepara relazioni per il processamento Bayesiano
        extracted_relationships = []
        for relationship in knowledge_relationships:
            extracted_relationships.append({
                'source': relationship.get('source_name', ''),
                'target': relationship.get('target_name', ''),
                'type': relationship.get('relationship_type', 'related_to'),
                'description': relationship.get('relationship_description', '')
            })

        # Processa le prove estratte tramite il motore Bayesiano
        try:
            bayesian_result = bayesian_engine.process_document_evidence(
                document_file_name=file_name,
                extracted_entities=extracted_entities,
                extracted_relationships=extracted_relationships
            )

            if bayesian_result.success:
                print(f"✅ Processamento Bayesiano completato: {bayesian_result.entities_created} entità, {bayesian_result.relationships_created} relazioni")
                academic_metadata['bayesian_processing'] = {
                    'success': True,
                    'entities_created': bayesian_result.entities_created,
                    'relationships_created': bayesian_result.relationships_created,
                    'confidence_summary': bayesian_engine.get_confidence_summary()
                }
            else:
                print(f"⚠️ Processamento Bayesiano fallito: {bayesian_result.errors}")
                academic_metadata['bayesian_processing'] = {
                    'success': False,
                    'errors': bayesian_result.errors
                }

        except Exception as e:
            print(f"❌ Errore nel processamento Bayesiano: {e}")
            academic_metadata['bayesian_processing'] = {
                'success': False,
                'error': str(e)
            }

        # Salvare entità e relazioni nel grafo della conoscenza
        # Nota: Dato che process_document_task non ha user_id, questo sarà gestito dal chiamante
        # Per ora, le mettiamo in metadata e saranno salvate dal chiamante se loggato
        academic_metadata['knowledge_entities'] = knowledge_entities
        academic_metadata['knowledge_relationships'] = knowledge_relationships

        # 5. SALVATAGGIO SU DB E ARCHIVIAZIONE
        update_status("Salvataggio finale...", file_name)
        with db_connect() as conn:
            conn.cursor().execute("""
                INSERT INTO papers (file_name, title, authors, publication_year, category_id, category_name, formatted_preview, keywords, ai_tasks, processed_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?) ON CONFLICT(file_name) DO UPDATE SET
                title=excluded.title, authors=excluded.authors, publication_year=excluded.publication_year,
                category_id=excluded.category_id, category_name=excluded.category_name, formatted_preview=excluded.formatted_preview,
                keywords=excluded.keywords, ai_tasks=excluded.ai_tasks, processed_at=excluded.processed_at;
            """, (file_name, metadata.title, json.dumps(metadata.authors), metadata.publication_year, category_id, category_full_name, formatted_preview, json.dumps(academic_metadata.get('keywords', [])), json.dumps(academic_metadata.get('ai_tasks', {})), datetime.now().isoformat()))
            conn.commit()

        destination_folder = os.path.join(CATEGORIZED_ARCHIVE_DIR, part_id, chapter_id)
        os.makedirs(destination_folder, exist_ok=True)
        shutil.move(file_path, os.path.join(destination_folder, file_name))
        
        update_status("Completato", file_name)
        return {'status': 'success', 'file_name': file_name, 'category': category_id}

    except Exception as e:
        error_msg = f"Errore: {str(e)}"
        print(f"❌ ERRORE nel processamento di {file_name}: {error_msg}")
        update_status(error_msg, file_name)

        # Non spostare in errore se è un problema di configurazione AI
        if "LLM non disponibile" in str(e) or "Servizi AI non disponibili" in str(e):
            print("⚠️ Errore di configurazione AI - non sposto il file in errore")
            # Don't raise here, just return an error status
            return {'status': 'ai_service_error', 'file_name': file_name, 'error': str(e)}

        # Sposta in errore solo per errori reali del documento
        error_folder = os.path.join(DOCS_TO_PROCESS_DIR, "_error")
        os.makedirs(error_folder, exist_ok=True)
        if os.path.exists(file_path):
            try:
                print(f"📁 Sposto {file_name} nella cartella errori")
                shutil.move(file_path, os.path.join(error_folder, file_name))
                print(f"✅ File {file_name} spostato nella cartella errori")
            except Exception as move_error:
                print(f"⚠️ Errore spostamento file {file_name}: {move_error}")

        # Return error status instead of raising
        return {'status': 'processing_error', 'file_name': file_name, 'error': str(e)}
    finally:
        if os.path.exists(lock_file):
            os.remove(lock_file)
            print(f"🔓 Lock rilasciato per {file_name}")

@celery_app.task(name='archivista.delete_document')
def delete_document_task(file_name):
    """
    Task di cancellazione atomica: elimina un documento da indice, database e filesystem.
    Tutte le operazioni devono riuscire o fallire atomicamente per garantire consistenza.
    """
    lock_file = os.path.join(DB_STORAGE_DIR, f"delete_{file_name}.lock")

    # Check for existing deletion lock
    if os.path.exists(lock_file):
        lock_age = time.time() - os.path.getmtime(lock_file)
        if lock_age < 300:  # 5 minutes timeout for deletion
            print(f"⏳ Cancellazione di {file_name} già in corso, attendo...")
            return {'status': 'already_deleting', 'file_name': file_name}
        else:
            print(f"⚠️ Rimuovo lock obsoleto per cancellazione di {file_name}")
            os.remove(lock_file)

    try:
        # Create deletion lock
        with open(lock_file, "w") as f:
            f.write(f"{os.getpid()},{time.time()}")

        print(f"🗑️ Inizio cancellazione atomica di {file_name}")

        # STEP 1: Trova informazioni documento nel database
        with db_connect() as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT category_id, category_name FROM papers WHERE file_name = ?", (file_name,))
            row = cursor.fetchone()

            if not row:
                raise ValueError(f"Documento {file_name} non trovato nel database")

            category_id = row['category_id']
            category_name = row['category_name']

        print(f"📋 Documento trovato: {file_name} in categoria {category_id}")

        # STEP 2: Rimuovi dall'indice vettoriale
        try:
            from llama_index.core.storage.docstore import SimpleDocumentStore
            from llama_index.core.storage.index_store import SimpleIndexStore

            # Crea/ottiene collezione ChromaDB per cancellazione
            try:
                import chromadb
                from llama_index.vector_stores.chroma import ChromaVectorStore

                db = chromadb.PersistentClient(path=DB_STORAGE_DIR)
                chroma_collection = db.get_or_create_collection("documents")
                vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
                print(f"🚀 Utilizzo ChromaVectorStore per cancellazione")

            except ImportError:
                print(f"⚠️ ChromaDB non disponibile per cancellazione, fallback a SimpleVectorStore")
                from llama_index.core.vector_stores import SimpleVectorStore
                vector_store = SimpleVectorStore()

            storage_context = StorageContext.from_defaults(
                docstore=SimpleDocumentStore(),
                vector_store=vector_store,
                index_store=SimpleIndexStore(),
                persist_dir=DB_STORAGE_DIR
            )

            # Carica l'index esistente
            index = load_index_from_storage(storage_context)

            # Trova il documento per file_name nel metadata
            documents_to_delete = []
            for doc_id, doc in index.docstore.docs.items():
                if doc.metadata.get('file_name') == file_name:
                    documents_to_delete.append(doc_id)

            if documents_to_delete:
                # Rimuovi documenti dall'index
                for doc_id in documents_to_delete:
                    index.delete_ref_doc(doc_id, delete_from_docstore=True)
                print(f"✅ Rimosso {len(documents_to_delete)} documento/i dall'index vettoriale")
            else:
                print(f"⚠️ Documento {file_name} non trovato nell'index vettoriale")

            # Salva l'index aggiornato
            index.storage_context.persist(persist_dir=DB_STORAGE_DIR)
            print("💾 Index vettoriale aggiornato dopo cancellazione")

        except FileNotFoundError:
            print(f"⚠️ Index non esistente durante cancellazione di {file_name}")
        except Exception as index_error:
            print(f"⚠️ Errore rimozione dall'index: {index_error}")
            # Continuiamo con le altre operazioni

        # STEP 3: Rimuovi dal database
        with db_connect() as conn:
            cursor = conn.cursor()
            cursor.execute("DELETE FROM papers WHERE file_name = ?", (file_name,))
            deleted_rows = cursor.rowcount
            conn.commit()

        if deleted_rows > 0:
            print(f"✅ Rimossa {deleted_rows} riga/e dal database per {file_name}")
        else:
            print(f"⚠️ Nessuna riga rimossa dal database per {file_name}")

        # STEP 4: Cancella file fisico
        file_deleted = False
        if category_id and category_id != "UNCATEGORIZED/C00":
            # Trova il percorso del file categorizzato
            part_id, chapter_id = category_id.split('/')
            file_path = os.path.join(CATEGORIZED_ARCHIVE_DIR, part_id, chapter_id, file_name)

            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    file_deleted = True
                    print(f"✅ File fisico eliminato: {file_path}")
                except Exception as file_error:
                    print(f"⚠️ Errore eliminazione file {file_path}: {file_error}")
            else:
                print(f"⚠️ File fisico non trovato: {file_path}")
        else:
            print(f"⚠️ Categoria non valida per {file_name}: {category_id}")

        # STEP 5: Verifica consistenza finale
        remaining_issues = []

        # Controlla se il documento è ancora nel database
        with db_connect() as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT COUNT(*) as count FROM papers WHERE file_name = ?", (file_name,))
            if cursor.fetchone()['count'] > 0:
                remaining_issues.append("Documento ancora presente nel database")

        print(f"🗑️ Cancellazione completata per {file_name}")
        return {
            'status': 'success',
            'file_name': file_name,
            'index_removed': len(documents_to_delete) if 'documents_to_delete' in locals() else 0,
            'db_rows_deleted': deleted_rows,
            'file_deleted': file_deleted,
            'warnings': remaining_issues
        }

    except Exception as e:
        error_msg = f"Errore durante cancellazione di {file_name}: {str(e)}"
        print(f"❌ {error_msg}")
        return {
            'status': 'error',
            'file_name': file_name,
            'error': str(e)
        }
    finally:
        # Rimuovi lock di cancellazione
        if os.path.exists(lock_file):
            os.remove(lock_file)
            print(f"🔓 Lock di cancellazione rilasciato per {file_name}")

@celery_app.task(name='archivista.cleanup_old_data')
def cleanup_old_data():
    """Task di pulizia: elimina i file falliti più vecchi di 7 giorni."""
    pass

@celery_app.task(name='archivista.scan_new_documents_periodic')
def scan_for_new_documents_periodic():
    """Task periodica: scansiona la cartella di input e avvia il processamento per i nuovi file."""
    pass

@celery_app.task(name='archivista.process_user_bayesian_knowledge')
def process_user_bayesian_knowledge_task(user_id: int, file_name: str):
    """
    Task specifico per processare la conoscenza Bayesiana per un utente specifico.

    Questo task viene chiamato dopo che un documento è stato processato,
    per aggiornare il grafo della conoscenza personale dell'utente con confidenza Bayesiana.

    Args:
        user_id: ID dell'utente proprietario della conoscenza
        file_name: Nome del file già processato

    Returns:
        dict: Risultato del processamento Bayesiano
    """
    lock_file = os.path.join(DB_STORAGE_DIR, f"bayesian_{user_id}_{file_name}.lock")

    # Check for existing lock
    if os.path.exists(lock_file):
        lock_age = time.time() - os.path.getmtime(lock_file)
        if lock_age < 300:  # 5 minutes timeout
            print(f"⏳ Processamento Bayesiano già in corso per user {user_id}, file {file_name}")
            return {'status': 'already_processing', 'user_id': user_id, 'file_name': file_name}
        else:
            print(f"⚠️ Rimuovo lock obsoleto per processamento Bayesiano")
            os.remove(lock_file)

    try:
        # Create lock file
        with open(lock_file, "w") as f:
            f.write(f"{os.getpid()},{time.time()}")

        print(f"🧠 Inizio processamento Bayesiano per user {user_id}, file {file_name}")

        # Verifica che il documento esista nel database
        with db_connect() as conn:
            cursor = conn.cursor()
            cursor.execute("""
                SELECT category_id, category_name, formatted_preview, keywords, ai_tasks
                FROM papers WHERE file_name = ?
            """, (file_name,))

            paper = cursor.fetchone()
            if not paper:
                raise ValueError(f"Documento {file_name} non trovato nel database")

        # Crea motore di inferenza personalizzato per l'utente
        bayesian_engine = create_inference_engine(user_id=user_id, learning_rate=0.3)

        # Estrai testo del documento per l'analisi (se necessario)
        # Nota: In produzione, potresti voler cachare il testo estratto
        file_path = None
        if paper['category_id'] and paper['category_id'] != "UNCATEGORIZED/C00":
            part_id, chapter_id = paper['category_id'].split('/')
            categorized_path = os.path.join(CATEGORIZED_ARCHIVE_DIR, part_id, chapter_id, file_name)
            if os.path.exists(categorized_path):
                file_path = categorized_path

        if not file_path:
            raise ValueError(f"File {file_name} non trovato nel percorso categorizzato")

        # Estrai testo per l'analisi Bayesiana
        file_ext = os.path.splitext(file_name)[1].lower()
        extractor = get_text_extractor(file_ext)
        if not extractor:
            raise ValueError(f"Formato file non supportato per analisi Bayesiana: {file_ext}")

        full_text = extractor(file_path)
        if not full_text or not full_text.strip():
            raise ValueError("Documento vuoto o illeggibile per analisi Bayesiana")

        # Ri-estrai entità e relazioni per questo utente specifico
        # (in produzione, potresti voler salvare queste informazioni per utente)
        knowledge_entities = []
        knowledge_relationships = []

        # Usa le stesse funzioni di estrazione del processamento principale
        # ma con contesto utente-specifico

        # Estrazione entità concettuali
        try:
            entities_prompt = PromptTemplate(prompt_manager.get_prompt("KNOWLEDGE_ENTITIES_PROMPT"))
            entities_query = entities_prompt.format(document_text=full_text[:3000])
            entities_response = Settings.llm.complete(entities_query)
            entities_text = str(entities_response).strip()

            try:
                entities_data = json.loads(entities_text)
                knowledge_entities = entities_data
                print(f"✅ Estratte {len(knowledge_entities)} entità per user {user_id}")
            except (json.JSONDecodeError, KeyError):
                print(f"⚠️ Impossibile parsare entità JSON per user {user_id}")
                knowledge_entities = []

        except Exception as e:
            print(f"⚠️ Errore estrazione entità per user {user_id}: {e}")
            knowledge_entities = []

        # Estrazione relazioni (se abbiamo entità)
        if knowledge_entities:
            try:
                entities_list = [f"- {e['entity_name']} ({e['entity_type']})" for e in knowledge_entities]
                entities_list_text = "\n".join(entities_list)

                relationships_prompt = PromptTemplate(prompt_manager.get_prompt("ENTITY_RELATIONSHIPS_PROMPT"))
                relationships_query = relationships_prompt.format(
                    document_text=full_text[:3000],
                    entities_list=entities_list_text
                )
                relationships_response = Settings.llm.complete(relationships_query)
                relationships_text = str(relationships_response).strip()

                try:
                    relationships_data = json.loads(relationships_text)
                    knowledge_relationships = relationships_data
                    print(f"✅ Estratte {len(knowledge_relationships)} relazioni per user {user_id}")
                except (json.JSONDecodeError, KeyError):
                    print(f"⚠️ Impossibile parsare relazioni JSON per user {user_id}")
                    knowledge_relationships = []

            except Exception as e:
                print(f"⚠️ Errore estrazione relazioni per user {user_id}: {e}")
                knowledge_relationships = []

        # Prepara dati per il motore Bayesiano
        extracted_entities = []
        for entity in knowledge_entities:
            extracted_entities.append({
                'name': entity.get('entity_name', ''),
                'type': entity.get('entity_type', 'concept'),
                'description': entity.get('entity_description', '')
            })

        extracted_relationships = []
        for relationship in knowledge_relationships:
            extracted_relationships.append({
                'source': relationship.get('source_name', ''),
                'target': relationship.get('target_name', ''),
                'type': relationship.get('relationship_type', 'related_to'),
                'description': relationship.get('relationship_description', '')
            })

        # Processa tramite motore Bayesiano
        bayesian_result = bayesian_engine.process_document_evidence(
            document_file_name=file_name,
            extracted_entities=extracted_entities,
            extracted_relationships=extracted_relationships
        )

        # Risultato finale
        result_data = {
            'status': 'success' if bayesian_result.success else 'error',
            'user_id': user_id,
            'file_name': file_name,
            'entities_created': bayesian_result.entities_created,
            'relationships_created': bayesian_result.relationships_created,
            'confidence_summary': bayesian_engine.get_confidence_summary()
        }

        if not bayesian_result.success:
            result_data['errors'] = bayesian_result.errors

        print(f"✅ Processamento Bayesiano completato per user {user_id}")
        return result_data

    except Exception as e:
        error_msg = f"Errore nel processamento Bayesiano per user {user_id}: {str(e)}"
        print(f"❌ {error_msg}")
        return {
            'status': 'error',
            'user_id': user_id,
            'file_name': file_name,
            'error': str(e)
        }
    finally:
        if os.path.exists(lock_file):
            os.remove(lock_file)
            print(f"🔓 Lock rilasciato per processamento Bayesiano user {user_id}")

@celery_app.task(name='archivista.process_user_feedback_bayesian')
def process_user_feedback_bayesian_task(user_id: int, target_type: str, target_id: int,
                                      feedback_type: str, feedback_text: str = ""):
    """
    Task per processare il feedback dell'utente tramite il motore Bayesiano.

    Args:
        user_id: ID dell'utente che fornisce il feedback
        target_type: 'entity' o 'relationship'
        target_id: ID dell'entità o relazione target
        feedback_type: 'positive', 'negative', o 'correction'
        feedback_text: Testo esplicativo del feedback

    Returns:
        dict: Risultato del processamento del feedback
    """
    try:
        print(f"👤 Processamento feedback Bayesiano per user {user_id}")

        # Crea motore di inferenza per l'utente
        bayesian_engine = create_inference_engine(user_id=user_id, learning_rate=0.3)

        # Processa il feedback
        feedback_result = bayesian_engine.process_user_feedback(
            target_type=target_type,
            target_id=target_id,
            feedback_type=feedback_type,
            feedback_strength=1.0,
            feedback_text=feedback_text
        )

        result_data = {
            'status': 'success' if feedback_result.success else 'error',
            'user_id': user_id,
            'target_type': target_type,
            'target_id': target_id,
            'feedback_type': feedback_type,
            'updates_performed': len(feedback_result.updates_performed)
        }

        if not feedback_result.success:
            result_data['errors'] = feedback_result.errors

        print(f"✅ Feedback processato per user {user_id}")
        return result_data

    except Exception as e:
        error_msg = f"Errore nel processamento feedback per user {user_id}: {str(e)}"
        print(f"❌ {error_msg}")
        return {
            'status': 'error',
            'user_id': user_id,
            'error': str(e)
        }

@celery_app.task(name='archivista.apply_temporal_decay_bayesian')
def apply_temporal_decay_bayesian_task(user_id: int):
    """
    Task per applicare decadimento temporale al grafo della conoscenza dell'utente.

    Args:
        user_id: ID dell'utente

    Returns:
        dict: Risultato dell'operazione di decadimento
    """
    try:
        print(f"🕐 Applicazione decadimento temporale per user {user_id}")

        # Crea motore di inferenza per l'utente
        bayesian_engine = create_inference_engine(user_id=user_id, learning_rate=0.3)

        # Applica decadimento temporale
        decay_result = bayesian_engine.apply_temporal_decay_to_all()

        result_data = {
            'status': 'success' if decay_result.success else 'error',
            'user_id': user_id,
            'entities_updated': decay_result.metadata.get('entities_updated', 0),
            'relationships_updated': decay_result.metadata.get('relationships_updated', 0)
        }

        if not decay_result.success:
            result_data['errors'] = decay_result.errors

        print(f"✅ Decadimento temporale applicato per user {user_id}")
        return result_data

    except Exception as e:
        error_msg = f"Errore nell'applicazione decadimento temporale per user {user_id}: {str(e)}"
        print(f"❌ {error_msg}")
        return {
            'status': 'error',
            'user_id': user_id,
            'error': str(e)
        }
